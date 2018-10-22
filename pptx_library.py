# -*- coding: utf-8 -*-
"""
Created on Fri Aug 31 15:08:40 2018

@author: t58830
"""
import pptx
from pptx.util import Inches, Cm, Pt

class python_pptx_funcs(object):  
  """
  This function deletes a slide based on the index of the slide.
  
  All indices are zero-based.
  
  Arguments:
    * presentation - The presentation object (object)
    * index        - Index of the slide (integer)
  
  """
  def delete_slide(presentation, index):
    xml_slides = presentation.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[index])    
  
  """
  This function allows you to quickly locate a paragraph object on a template
  by looking for a keyword. Useful for textboxes and similar object.
  
  
  Arguments: 
    *  - slide       - Slide object (object)
    *  - target_text - Keyword to locate (string)
  """                  
  def locate_paragraph(slide, target_text):
    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          if run.text == target_text:
            return paragraph

  """
  This function allows you to quickly locate a shape object on a template
  by looking for a keyword. Useful for when locate_paragraph fails (edge cases) 
  
  Arguments: 
    *  - slide       - Slide object (object)
    *  - target_text - Keyword to locate (string)
  """              
  def locate_shape(slide, target_text):
    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          if run.text == target_text:
            return shape.text_frame
  
  """
  This function allows you replace the text on template on a specified slide
  with text of your choice. This text can be bold and/or italic.
  
  All indices are zero-based.  
  
  Arguments: 
    * slide       - The slide object (object)
    * input_text  - Text which you want to write (string)
    * target_text - The text of the run which you are trying to replace 
                    (string) 
    * font        - Optional argument, Specify which font (string)
    * font_size   - Optional argument, Size of the font (Integer)
    * bold        - Optional argument, pass True for bold text (boolean)
    * italic      - Optional argument, pass True for italic text (boolean)
  """  
  def run_text(slide,
               input_text,
               target_text,
               font = None,
               font_size = None,
               bold   = False,
               italic = False):
    run = ""

    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for paragraph in shape.text_frame.paragraphs:
        for iterrun in paragraph.runs:
          if iterrun.text == target_text:
            run = iterrun
    if run == "":
      print("Target could be not found: " + target_text)
            
    run.text = input_text
    if bold:
      run.font.bold = True
    if italic:
      run.font.italic = True
    if font != None:
      run.font.name = font
    if font_size != None:
      run.font.size = Pt(font_size)   
    
  """
  This

  All indices are zero-based.  
  
  Arguments:
    * table   - The table object (object)
    * row_idx - The row index    (integer)
    * col_idx - The column index (integr())
    * font    - Optional variable which allows you to pass a font name (string)
    * bold    - Optional argument, pass True for bold text (boolean)
    * italic  - Optional argument, pass True for italic text (boolean)
  """  
  def bold_table_cell(table,
                      row_idx,
                      col_idx,
                      font   = None,
                      bold   = False,
                      italic = False):
                        
    cell = table.cell(row_idx,col_idx)
    for paragraph in cell.text_frame.paragraphs:
      for run in paragraph.runs:
        if font is not None:
          run.font.name = font
        if bold:
          run.font.bold = True
        if italic:
          run.font.italic = True
        run.font.color.rgb = pptx.dml.color.RGBColor(0,0,0)
        
        
  """
  This function allows you to quickly create a table using a Pandas DatFrame
  on a slide of size which you have to specify. The table is created, 
  dynamically based on the size of DataFrame.
  
  Arguments:
    * slide  - The slide object (object)
    * data   - Pandas DataFrame
    * left   - x postion of where the table should start (in Cm) (integer)
    * top    - y position of where the table should end (in Cm) (integer)
    * width  - x width (in Cm) (integer)
    * height - y height (in Cm) (integer)
    * header - Optional variable which determines whether the header of the 
               Pandas DataFrame is included in the table. Send False to not
               include. (boolean)
  """
  def create_table(slide,
                   data,
                   left,
                   top,
                   width,
                   height,
                   header = True):
                     
    if header:   
      data.loc[-1] = list(data.columns)  # adding a row
      data.index = data.index + 1  # shifting index
      data.sort_index(inplace=True)
    
    left   = Cm(left)
    top    = Cm(top)
    width  = Cm(width)
    height = Cm(height)  

    n_rows = len(data)
    n_cols = len(data.columns)            
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for i, column in enumerate(data.columns):    
      for j in range(len(data[column])):
        table.cell(j, i).text = str(data[column][j])   
        table.cell(j, i).vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
    return table
    
  """
  This functions allows you to easily merge two or more cells. It keeps the
  content in the leftmost cell.
  
  All indices are zero-based.
  
  Arguments:
    * table         - The table object
    * start_col_idx - Start column of the cells you want to merge
    * end_col_idx   - End column of the cells you want to merge
    * row_idx       - The row in which you merge cells
  """
  def mergeCellsHorizontally(table,
                             start_col_idx,
                             end_col_idx,
                             row_idx):
                               
    col_count = end_col_idx - start_col_idx + 1
    row_cells = [c for c in table.rows[row_idx].cells][start_col_idx:end_col_idx]
    row_cells[0]._tc.set('gridSpan', str(col_count))
  
  """
  This is a generator which is extremely useful if you want to apply a setting
  to the whole table. E.G. setting a font to all cells on the table.
  
  It can be used in the below way (pseudocode)
    
  _____________________________________________________________________________ 
    for cell in iter_cells(table):
      do something with cell
  _____________________________________________________________________________
  
  Arguments:
    * table - The table object
  """
  def iter_cells(table):
    for row in table.rows:
       for cell in row.cells:
          yield cell
          
  """
  This is useful if you want to move slides around after preesentation has been
  created.
  
  All indices are zero-based.
   
  Arguments:
    presentation - The presentation object
    old_idx    - The old index of the slide you want to move
    new_idx    - The new index of the slide you want to move
    
  """          
  def move_slide(presentation, old_idx, new_idx):
      xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
      slides = list(xml_slides)
      xml_slides.remove(slides[old_idx])
      xml_slides.insert(new_idx, slides[old_idx])

#The following functions still need to be documented

  def inject_text(self, slide, text, injection_list):
    for text, target in zip(text_list, injection_list):
      python_pptx_funcs.run_text(self.slide,text,target)

    